from io import BytesIO
from datetime import datetime

import pandas as pd
from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from Backend.scrapers.bcra_scraper import get_bcra_bancos

router = APIRouter()

GREEN = "86BC25"
LBLUE = "DCE6F1"
GRAY = "F5F5F5"
WHITE = "FFFFFF"

thin = Border(
    left=Side(style="thin"),
    right=Side(style="thin"),
    top=Side(style="thin"),
    bottom=Side(style="thin"),
)

AL = Alignment(horizontal="left", vertical="center")
AR = Alignment(horizontal="right", vertical="center")
AC = Alignment(horizontal="center", vertical="center")


def hdr_font():
    return Font(name="Arial", bold=True, color="FFFFFF", size=9)


def dat_font():
    return Font(name="Arial", size=9)


def bold_font():
    return Font(name="Arial", bold=True, size=9)


def title_font():
    return Font(name="Arial", bold=True, size=13, color=GREEN)


def sub_font():
    return Font(name="Arial", size=9, italic=True, color="666666")


def hdr_fill():
    return PatternFill("solid", start_color=GREEN)


def sum_fill():
    return PatternFill("solid", start_color=LBLUE)


def alt_fill(i):
    return PatternFill("solid", start_color=GRAY if i % 2 == 0 else WHITE)


def num(v):
    if v is None or v == "":
        return None
    try:
        return float(v)
    except:
        return None


def build_bcra_dataframe():

    data = get_bcra_bancos()

    bancos = data.get("bancos", [])

    fecha_reporte = data.get("fecha_reporte", "N/D")
    fecha_export = datetime.now().strftime("%d/%m/%Y")

    rows = []

    for b in bancos:

        activos = num(b.get("Activos"))
        depositos = num(b.get("Depositos"))
        prestamos = num(b.get("Prestamos"))
        patrimonio = num(b.get("Patrimonio Neto"))

        rows.append(
            {
                "Banco": b.get("Banco", "N/A"),
                "Activos": activos,
                "Depósitos": depositos,
                "Préstamos": prestamos,
                "Patrimonio Neto": patrimonio,
            }
        )

    df = pd.DataFrame(rows)

    if df.empty:
        return df, fecha_reporte, fecha_export

    df = df.sort_values("Activos", ascending=False).reset_index(drop=True)

    return df, fecha_reporte, fecha_export


def write_header_row(ws, row, cols):

    for ci, (name, width) in enumerate(cols, 1):

        cell = ws.cell(row=row, column=ci, value=name)

        cell.font = hdr_font()
        cell.fill = hdr_fill()
        cell.alignment = AC
        cell.border = thin

        ws.column_dimensions[get_column_letter(ci)].width = width


def write_data_row(ws, row_data, cols, excel_row):

    fill = alt_fill(excel_row)

    for ci, key in enumerate(cols, 1):

        val = row_data.get(key)

        cell = ws.cell(row=excel_row, column=ci, value=val)

        cell.font = dat_font()
        cell.fill = fill
        cell.border = thin

        if ci > 1:
            cell.alignment = AR
            cell.number_format = '#,##0.0'
        else:
            cell.alignment = AL


def add_summary_rows(ws, data_start, last_row):

    summary_start = last_row + 2

    ws.cell(summary_start, 1, "Mediana").font = bold_font()
    ws.cell(summary_start + 1, 1, "Promedio").font = bold_font()

    ws.cell(summary_start, 1).fill = sum_fill()
    ws.cell(summary_start + 1, 1).fill = sum_fill()

    for col in range(2, 6):

        col_letter = get_column_letter(col)

        med = ws.cell(
            summary_start,
            col,
            f"=MEDIAN({col_letter}{data_start}:{col_letter}{last_row})",
        )

        avg = ws.cell(
            summary_start + 1,
            col,
            f"=AVERAGE({col_letter}{data_start}:{col_letter}{last_row})",
        )

        for c in [med, avg]:
            c.font = bold_font()
            c.fill = sum_fill()
            c.border = thin
            c.alignment = AR
            c.number_format = '#,##0.0'


def generate_excel(df, fecha_reporte, fecha_export):

    wb = Workbook()

    cols = [
        ("Banco", 34),
        ("Activos", 18),
        ("Depósitos", 18),
        ("Préstamos", 18),
        ("Patrimonio Neto", 20),
    ]

    # Hoja 1
    ws = wb.active
    ws.title = "BCRA Ranking"

    ws.merge_cells("A1:E1")
    ws["A1"] = "BCRA Intelligence — Sistema Financiero Argentino"
    ws["A1"].font = title_font()

    ws.merge_cells("A2:E2")
    ws["A2"] = f"Reporte BCRA: {fecha_reporte} | Exportado: {fecha_export}"
    ws["A2"].font = sub_font()

    write_header_row(ws, 4, cols)

    start = 5

    for i, (_, r) in enumerate(df.iterrows(), start):

        write_data_row(
            ws,
            r.to_dict(),
            [c[0] for c in cols],
            i,
        )

    last_row = start + len(df) - 1

    add_summary_rows(ws, start, last_row)

    # Hoja 2
    ws2 = wb.create_sheet("Top 10")

    ws2.merge_cells("A1:E1")
    ws2["A1"] = "Top 10 por Activos"
    ws2["A1"].font = title_font()

    ws2.merge_cells("A2:E2")
    ws2["A2"] = f"Reporte BCRA: {fecha_reporte}"
    ws2["A2"].font = sub_font()

    write_header_row(ws2, 4, cols)

    df10 = df.head(10)

    start = 5

    for i, (_, r) in enumerate(df10.iterrows(), start):

        write_data_row(
            ws2,
            r.to_dict(),
            [c[0] for c in cols],
            i,
        )

    last_row = start + len(df10) - 1

    add_summary_rows(ws2, start, last_row)
   # ─────────────────────────────
    # HOJA 3 — DASHBOARD REAL
    # ─────────────────────────────
    from openpyxl.chart import BarChart, Reference, PieChart

    ws3 = wb.create_sheet("Dashboard")

    # TITULO
    ws3.merge_cells("A1:F1")
    ws3["A1"] = "BCRA Intelligence — Dashboard"
    ws3["A1"].font = title_font()

    # KPIs (tipo cards)
    kpi_labels = ["Activos", "Depósitos", "Préstamos", "Patrimonio"]

    for i, (label, col) in enumerate(zip(kpi_labels, ["B", "C", "D", "E"])):

        row = 3

        ws3.cell(row, i*2 + 1, label).font = bold_font()

        formula = f"=SUM('BCRA Ranking'!{col}5:{col}{4+len(df)})"

        val = ws3.cell(row+1, i*2 + 1, formula)
        val.number_format = '#,##0.0'
        val.font = Font(size=12, bold=True)

    # ─────────────────────────────
    # TABLA TOP 10 (para gráficos)
    # ─────────────────────────────
    ws3["A7"] = "Top 10 Bancos"
    ws3["A7"].font = bold_font()

    for i in range(10):
        ws3.cell(9+i, 1, f"='BCRA Ranking'!A{5+i}")
        ws3.cell(9+i, 2, f"='BCRA Ranking'!B{5+i}")

    # ─────────────────────────────
    # GRÁFICO BARRAS
    # ─────────────────────────────
    bar = BarChart()
    bar.title = "Top 10 por Activos"

    data = Reference(ws3, min_col=2, min_row=9, max_row=18)
    cats = Reference(ws3, min_col=1, min_row=9, max_row=18)

    bar.add_data(data)
    bar.set_categories(cats)

    bar.height = 8
    bar.width = 16

    ws3.add_chart(bar, "D7")

    # ─────────────────────────────
    # GRÁFICO PIE (Top 5)
    # ─────────────────────────────
    pie = PieChart()
    pie.title = "Market Share"

    data = Reference(ws3, min_col=2, min_row=9, max_row=13)
    cats = Reference(ws3, min_col=1, min_row=9, max_row=13)

    pie.add_data(data)
    pie.set_categories(cats)

    ws3.add_chart(pie, "D20")

    # tamaños columnas
    ws3.column_dimensions["A"].width = 30
    ws3.column_dimensions["B"].width = 20
    buffer = BytesIO()

    wb.save(buffer)

    buffer.seek(0)

    return buffer


@router.get("/bcra/export-excel")
def export_bcra_excel():

    df, fecha_reporte, fecha_export = build_bcra_dataframe()

    if df.empty:
        return {"error": "Sin datos para exportar"}

    buffer = generate_excel(df, fecha_reporte, fecha_export)

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={
            "Content-Disposition": "attachment; filename=bcra_ranking.xlsx"
        },
    )

from pptx import Presentation
from pptx.util import Inches, Pt


def generate_powerpoint(df, fecha_reporte, fecha_export):

    prs = Presentation()

    # Slide 1 - Título
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = "BCRA Intelligence"
    slide.placeholders[1].text = f"Sistema Financiero Argentino\nReporte: {fecha_reporte}"

    # Slide 2 - Tabla ranking
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "Ranking Bancos por Activos"

    rows = len(df) + 1
    cols = 5

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(5),
    ).table

    headers = ["Banco", "Activos", "Depósitos", "Préstamos", "Patrimonio"]

    for i, h in enumerate(headers):
        table.cell(0, i).text = h

    for r, (_, row) in enumerate(df.iterrows(), 1):

        table.cell(r, 0).text = str(row["Banco"])
        table.cell(r, 1).text = f'{row["Activos"]:,.1f}'
        table.cell(r, 2).text = f'{row["Depósitos"]:,.1f}'
        table.cell(r, 3).text = f'{row["Préstamos"]:,.1f}'
        table.cell(r, 4).text = f'{row["Patrimonio Neto"]:,.1f}'

    # Slide 3 - Top 10
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Top 10 Bancos"

    df10 = df.head(10)

    rows = len(df10) + 1

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(1.5),
        Inches(9),
        Inches(5),
    ).table

    for i, h in enumerate(headers):
        table.cell(0, i).text = h

    for r, (_, row) in enumerate(df10.iterrows(), 1):

        table.cell(r, 0).text = str(row["Banco"])
        table.cell(r, 1).text = f'{row["Activos"]:,.1f}'
        table.cell(r, 2).text = f'{row["Depósitos"]:,.1f}'
        table.cell(r, 3).text = f'{row["Préstamos"]:,.1f}'
        table.cell(r, 4).text = f'{row["Patrimonio Neto"]:,.1f}'

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer

@router.get("/bcra/export-ppt")
def export_bcra_ppt():

    df, fecha_reporte, fecha_export = build_bcra_dataframe()

    if df.empty:
        return {"error": "Sin datos"}

    buffer = generate_powerpoint(df, fecha_reporte, fecha_export)

    return StreamingResponse(
        buffer,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": "attachment; filename=bcra_ranking.pptx"
        },
    )